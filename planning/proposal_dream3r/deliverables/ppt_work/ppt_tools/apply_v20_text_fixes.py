from pathlib import Path

from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor


WORK = Path(r"E:\kykt\Dream\planning\proposal_dream3r\deliverables\ppt_work")
SRC = WORK / "proposal_dream3r_opening_report_reference_mode_v19.pptx"
OUT = WORK / "proposal_dream3r_opening_report_reference_mode_v20.pptx"


def shape_text(shape):
    return getattr(shape, "text", "")


def set_text_preserve_first_run(shape, text):
    """Replace visible text while keeping the first run's local formatting."""
    if not getattr(shape, "has_text_frame", False):
        return
    tf = shape.text_frame
    first_run = None
    for p in tf.paragraphs:
        for r in p.runs:
            if first_run is None:
                first_run = r
            r.text = ""
    if first_run is None:
        p = tf.paragraphs[0]
        first_run = p.add_run()
    first_run.text = text


def replace_exact(slide, old, new):
    count = 0
    for shape in slide.shapes:
        if shape_text(shape).strip() == old:
            set_text_preserve_first_run(shape, new)
            count += 1
    return count


def replace_contains(slide, old, new):
    count = 0
    for shape in slide.shapes:
        if not getattr(shape, "has_text_frame", False):
            continue
        for p in shape.text_frame.paragraphs:
            for r in p.runs:
                if old in r.text:
                    r.text = r.text.replace(old, new)
                    count += 1
    return count


def set_section_number(slide, num):
    candidates = []
    for shape in slide.shapes:
        txt = shape_text(shape).strip()
        if txt.isdigit() and shape.left < Inches(0.5) and shape.top < Inches(0.7):
            candidates.append(shape)
    if candidates:
        set_text_preserve_first_run(candidates[0], num)
        return True
    return False


def add_support_note(slide, x, y, w, text):
    box = slide.shapes.add_textbox(Inches(x), Inches(y), Inches(w), Inches(0.26))
    tf = box.text_frame
    tf.margin_left = 0
    tf.margin_right = 0
    tf.margin_top = 0
    tf.margin_bottom = 0
    p = tf.paragraphs[0]
    run = p.add_run()
    run.text = text
    run.font.name = "Microsoft YaHei"
    run.font.size = Pt(12)
    run.font.color.rgb = RGBColor(86, 104, 128)
    return box


def main():
    prs = Presentation(SRC)

    # 1. Fix section numbers.
    section_map = {
        3: "01", 4: "01", 5: "01",
        6: "02", 7: "02", 8: "02", 9: "02", 10: "02", 11: "02", 12: "02",
        13: "03", 14: "03", 15: "03", 16: "03", 17: "03", 18: "03",
        19: "04", 20: "04", 21: "04", 22: "04", 23: "04",
    }
    for slide_no, num in section_map.items():
        set_section_number(prs.slides[slide_no - 1], num)

    # 2. Claim simplification.
    claim_replacements = {
        3: ("3R 的价值在于统一输入输出；后续问题集中转向长序列、动态场景和结果可靠性。",
            "3R 统一输入输出，待解决长序列与可靠性"),
        4: ("Spann3R、CUT3R、Test3R、MASt3R 等工作提供了可借鉴的机制来源。",
            "四类机制为新架构提供直接参考"),
        5: ("课题目标由这些短板自然引出：记忆、校验、多专家编排和统一平台。",
            "五类短板引出四个设计方向"),
        6: ("本课题的架构由这些机制组合自然推出，并通过实验检验各模块作用。",
            "学习有效机制 + 引入新方法补足短板"),
        7: ("平台提供统一实验条件；新架构完成后可作为新模型接入平台进行同条件比较。",
            "架构是研究对象，平台是实验支撑"),
        13: ("聚合平台用于降低多模型实验的人工配置和手动对比成本。",
             "统一调度、统一格式、统一对比"),
        14: ("功能覆盖任务提交、状态追踪、样本对比、模型注册和结果归集。",
             "五项功能覆盖完整实验链路"),
        16: ("统一输入为 13 帧 1920×1080；耗时差异体现统一调度与记录的必要性。",
             "统一条件下耗时差异显著"),
        18: ("新架构训练完成后可作为新模型接入平台，与现有模型同条件比较。",
             "平台支撑架构验证，架构丰富平台模型池"),
    }
    for slide_no, (old, new) in claim_replacements.items():
        replace_exact(prs.slides[slide_no - 1], old, new)

    # 3. Global wording: use "multi-model" instead of "multi-expert" in this deck.
    for slide in prs.slides:
        replace_contains(slide, "多专家", "多模型")
        replace_contains(slide, "专家编排", "模型编排")
        replace_contains(slide, "专家切换", "模型切换")
        replace_contains(slide, "专家对照", "模型对照")
        replace_contains(slide, "专家各取所长", "模型各取所长")

    # 4. Slide 11: explicit label fixes.
    s11 = prs.slides[10]
    replace_exact(s11, "多模型编排模块", "多模型编排模块")
    replace_exact(s11, "多模型编排模块根据输入条件和校验结果在七个模型专家之间切换。",
                  "多模型编排模块根据输入条件和校验结果在七个候选模型之间切换。")
    replace_exact(s11, "能力匹配度跨度", "能力匹配")
    replace_exact(s11, "成本调整解析平局", "成本约束")
    replace_exact(s11, "失败退化", "失败回退")
    replace_exact(s11, "单模型各有优势场景，多模型路由用于提升复杂输入下的适应能力。",
                  "按场景、预算和校验结果选择合适模型")

    # 5. Slide 13: add reproducibility point without changing the row count.
    s13 = prs.slides[12]
    replace_exact(s13, "平台需求", "复现与平台需求")
    replace_exact(s13, "统一模型调度、结果格式化和多模型对照",
                  "依赖版本影响复现；需统一调度、格式和对比")

    # 6. Slide 19: add dataset plan.
    s19 = prs.slides[18]
    replace_contains(
        s19,
        "失败场景泛化：动态 / 弱纹理 / 大视角",
        "失败场景泛化：动态 / 弱纹理 / 大视角\n数据集：ScanNet / KITTI / TUM-RGBD",
    )

    # 7. Slide 20: wording and short support lines.
    s20 = prs.slides[19]
    replace_exact(s20, "三分支稀疏注意力记忆统一", "三分支稀疏注意力记忆模块")
    add_support_note(s20, 2.50, 3.02, 3.75, "支持在线检测与修复")
    add_support_note(s20, 8.19, 3.02, 3.75, "按场景按需调用")
    add_support_note(s20, 2.50, 4.58, 3.75, "覆盖远程/检索/局部时间尺度")
    add_support_note(s20, 8.19, 4.58, 3.75, "6 个模型接入，4 个端到端验证")

    prs.save(OUT)
    print(OUT)


if __name__ == "__main__":
    main()
