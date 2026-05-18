from pathlib import Path

from pptx import Presentation


WORK = Path(r"E:\kykt\Dream\planning\proposal_dream3r\deliverables\ppt_work")
SRC = WORK / "proposal_dream3r_opening_report_reference_mode_v20.pptx"
OUT = WORK / "proposal_dream3r_opening_report_reference_mode_v21.pptx"


def set_text_preserve_first_run(shape, text):
    if not getattr(shape, "has_text_frame", False):
        return
    tf = shape.text_frame
    first_run = None
    for p in tf.paragraphs:
        for r in p.runs:
            if first_run is None:
                first_run = r
            r.text = ""
    if first_run is None:
        first_run = tf.paragraphs[0].add_run()
    first_run.text = text


def main():
    prs = Presentation(SRC)
    cover = prs.slides[0]
    for shape in cover.shapes:
        if getattr(shape, "text", "").strip() == "XXX":
            set_text_preserve_first_run(shape, "崔昊喆 / 纪博闻")
            break
    prs.save(OUT)
    print(OUT)


if __name__ == "__main__":
    main()
