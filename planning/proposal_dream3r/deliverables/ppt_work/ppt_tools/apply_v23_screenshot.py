from pathlib import Path

from PIL import Image
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor


WORK = Path(r"E:\kykt\Dream\planning\proposal_dream3r\deliverables\ppt_work")
SRC = WORK / "proposal_dream3r_opening_report_reference_mode_v22_final.pptx"
OUT = WORK / "proposal_dream3r_opening_report_reference_mode_v23_with_screenshot.pptx"
SCREENSHOT = WORK / "软件截图.png"
CROPPED = WORK / "reference_mode_assets" / "software_screenshot_cropped.png"


def crop_screenshot():
    CROPPED.parent.mkdir(parents=True, exist_ok=True)
    im = Image.open(SCREENSHOT).convert("RGB")
    # Remove only the OS title bar so the slide focuses on the KYKT interface.
    cropped = im.crop((0, 42, im.width, im.height))
    cropped.save(CROPPED, quality=95)


def add_label(slide, x, y, text):
    box = slide.shapes.add_textbox(Inches(x), Inches(y), Inches(2.8), Inches(0.28))
    tf = box.text_frame
    tf.margin_left = 0
    tf.margin_right = 0
    tf.margin_top = 0
    tf.margin_bottom = 0
    run = tf.paragraphs[0].add_run()
    run.text = text
    run.font.name = "Microsoft YaHei"
    run.font.size = Pt(13)
    run.font.bold = True
    run.font.color.rgb = RGBColor(0, 70, 142)


def main():
    crop_screenshot()
    prs = Presentation(SRC)
    slide = prs.slides[14]  # Slide 15

    # Existing architecture figure.
    arch_pic = None
    for shape in slide.shapes:
        if shape.shape_type == 13:  # PICTURE
            arch_pic = shape
            break
    if arch_pic is not None:
        arch_pic.left = Inches(0.65)
        arch_pic.top = Inches(2.25)
        arch_pic.width = Inches(5.55)
        arch_pic.height = Inches(3.12)

    add_label(slide, 0.65, 1.93, "平台执行架构")
    add_label(slide, 7.15, 1.93, "KYKT Vision 运行界面")

    slide.shapes.add_picture(str(CROPPED), Inches(7.15), Inches(2.25), width=Inches(5.55))

    prs.save(OUT)
    print(OUT)
    print(CROPPED)


if __name__ == "__main__":
    main()
