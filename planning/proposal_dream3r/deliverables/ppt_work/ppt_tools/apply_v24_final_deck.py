from pathlib import Path

from PIL import Image
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor


WORK = Path(r"E:\kykt\Dream\planning\proposal_dream3r\deliverables\ppt_work")
SRC = WORK / "proposal_dream3r_opening_report_reference_mode_v23_with_screenshot.pptx"
OUT = WORK / "proposal_dream3r_opening_report_final.pptx"
COVER = WORK / "封面.png"
SHOT1 = WORK / "软件截图.png"
SHOT2 = WORK / "软件截图2.png"
ASSET_DIR = WORK / "reference_mode_assets"
SHOT1_CROP = ASSET_DIR / "software_screenshot_dashboard_crop.png"
SHOT2_CROP = ASSET_DIR / "software_screenshot_models_crop.png"


def remove_shape(shape):
    shape._element.getparent().remove(shape._element)


def crop_assets():
    ASSET_DIR.mkdir(parents=True, exist_ok=True)
    im1 = Image.open(SHOT1).convert("RGB")
    # Keep dashboard/task evidence and crop to the exact slide thumbnail ratio.
    im1.crop((0, 42, im1.width, min(im1.height, 666))).save(SHOT1_CROP, quality=95)
    im2 = Image.open(SHOT2).convert("RGB")
    # Keep model-route cards as a wide strip; avoids squeezing the screenshot in PPT.
    im2.crop((0, 42, im2.width, min(im2.height, 477))).save(SHOT2_CROP, quality=95)


def send_to_back(shape):
    sp_tree = shape._element.getparent()
    sp_tree.remove(shape._element)
    sp_tree.insert(2, shape._element)


def add_label(slide, x, y, w, text):
    box = slide.shapes.add_textbox(Inches(x), Inches(y), Inches(w), Inches(0.26))
    tf = box.text_frame
    tf.margin_left = 0
    tf.margin_right = 0
    tf.margin_top = 0
    tf.margin_bottom = 0
    run = tf.paragraphs[0].add_run()
    run.text = text
    run.font.name = "Microsoft YaHei"
    run.font.size = Pt(12)
    run.font.bold = True
    run.font.color.rgb = RGBColor(0, 70, 142)


def set_cover_background(prs):
    for slide_index in [0, 23]:
        slide = prs.slides[slide_index]
        # Remove the old grass/factory background and tinted wash; keep text, ribbons, and footers.
        for shape in list(slide.shapes):
            if shape.shape_type == 13 and shape.left <= Inches(0.05) and shape.width >= Inches(13):
                remove_shape(shape)
            elif shape.shape_type == 1 and shape.left <= Inches(0.02) and shape.top <= 0 and shape.width >= Inches(13) and shape.height >= Inches(5):
                remove_shape(shape)
        bg = slide.shapes.add_picture(str(COVER), 0, 0, width=prs.slide_width, height=prs.slide_height)
        send_to_back(bg)


def simplify_slide15(prs):
    slide = prs.slides[14]
    # Remove labels and screenshots previously inserted by v23.
    for shape in list(slide.shapes):
        txt = getattr(shape, "text", "").strip()
        if txt in {"平台执行架构", "KYKT Vision 运行界面"}:
            remove_shape(shape)
        elif shape.shape_type == 13 and shape.left >= Inches(7):
            remove_shape(shape)

    arch_pic = None
    for shape in slide.shapes:
        if shape.shape_type == 13:
            arch_pic = shape
            break

    if arch_pic is not None:
        arch_pic.left = Inches(0.62)
        arch_pic.top = Inches(2.20)
        arch_pic.width = Inches(5.55)
        arch_pic.height = Inches(3.12)
    add_label(slide, 0.62, 1.91, 2.4, "平台执行架构")

    add_label(slide, 6.75, 1.91, 3.0, "KYKT Vision：任务面板")
    slide.shapes.add_picture(str(SHOT1_CROP), Inches(6.75), Inches(2.18), width=Inches(5.95), height=Inches(2.08))

    add_label(slide, 6.75, 4.48, 3.2, "KYKT Vision：模型路线")
    slide.shapes.add_picture(str(SHOT2_CROP), Inches(6.75), Inches(4.75), width=Inches(5.95), height=Inches(1.45))

    # Keep the footer claim, but make room visually by lifting it a little if needed.
    for shape in slide.shapes:
        txt = getattr(shape, "text", "").strip()
        if txt.startswith("执行器只封装模型差异"):
            shape.top = Inches(6.42)
            shape.height = Inches(0.42)


def main():
    crop_assets()
    prs = Presentation(SRC)
    set_cover_background(prs)
    simplify_slide15(prs)
    prs.save(OUT)
    print(OUT)
    print(SHOT1_CROP)
    print(SHOT2_CROP)


if __name__ == "__main__":
    main()
