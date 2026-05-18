from pathlib import Path

from pptx import Presentation


WORK = Path(r"E:\kykt\Dream\planning\proposal_dream3r\deliverables\ppt_work")
SRC = WORK / "proposal_dream3r_opening_report_reference_mode_v21.pptx"
OUT = WORK / "proposal_dream3r_opening_report_reference_mode_v22_final.pptx"


def set_text_preserve_first_run(shape, text):
    if not getattr(shape, "has_text_frame", False):
        return False
    tf = shape.text_frame
    first_run = None
    for p in tf.paragraphs:
        for r in p.runs:
            if first_run is None:
                first_run = r
            r.text = ""
    if first_run is None:
        first_run = tf.paragraphs[0].add_run()
    first_run.text = text
    return True


def replace_exact(slide, old, new):
    count = 0
    for shape in slide.shapes:
        if getattr(shape, "text", "").strip() == old:
            if set_text_preserve_first_run(shape, new):
                count += 1
    return count


def replace_contains(slide, old, new):
    count = 0
    for shape in slide.shapes:
        if not getattr(shape, "has_text_frame", False):
            continue
        for p in shape.text_frame.paragraphs:
            for r in p.runs:
                if old in r.text:
                    r.text = r.text.replace(old, new)
                    count += 1
    return count


def main():
    prs = Presentation(SRC)

    replacements = {
        4: [
            ("3R 模型的优势集中在速度、端到端学习、统一输出和特色机制。",
             "3R 用一次推理替代多步级联，后续工作还沉淀出了四类可复用机制。"),
        ],
        5: [
            ("现有 3R 模型在长序列、动态场景、自检查和多模型协同方面仍有短板。",
             "速度和统一性解决了，但长序列、动态场景和结果可靠性问题依然突出。"),
        ],
        6: [
            ("设计思路分为两步：提取 3R 有效机制，再引入新方法补足短板。",
             "先从已有 3R 工作中提炼可用机制，再用新方法补上它们没解决的部分。"),
            ("学习有效机制 + 引入新方法补足短板",
             "六项机制 + 四种新方法 → 候选架构方案"),
        ],
        7: [
            ("课题包含新架构设计和聚合管理平台两条主线，二者服务于同一组实验目标。",
             "架构负责出方案，平台负责跑实验和做对比。"),
        ],
        8: [
            ("总体架构以总线为核心，把感知、记忆、动静分离、校验和模型编排串联起来。",
             "五个功能模块通过总线联动，输入图像序列后逐帧输出三维点图和校验证据。"),
        ],
        9: [
            ("空间记忆模块把长序列上下文拆成压缩状态、空间锚点和局部滑窗三条路径。",
             "长序列不能全靠全局注意力，需要把远程、检索和局部三种记忆分开处理。"),
        ],
        10: [
            ("校验理念来自 Test3R，本课题把它接入架构内部的修复动作。",
             "校验思路源自 Test3R，本课题将其嵌入架构内部，支持在线修复。"),
        ],
        13: [
            ("多模型对比实验需要统一调度、统一格式和统一结果归集。",
             "六个模型各有各的环境和格式，手动对齐效率太低、结果难复现。"),
        ],
        14: [
            ("平台围绕任务提交、状态追踪、结果对比、模型注册和远程调度展开。",
             "从任务提交到结果归档，平台覆盖多模型实验的完整链路。"),
        ],
        16: [
            ("6 个模型端到端验证通过，统一输入下的推理耗时差异明显。",
             "6 个模型端到端验证通过，同一输入下最快与最慢相差约 9 倍。"),
        ],
        17: [
            ("平台后续将围绕模型验证、对比视图、报告导出和下游接口推进。",
             "短期补齐验证，中期支撑论文，长期对接下游应用。"),
        ],
        19: [
            ("实验设计分为架构有效性验证和平台支撑能力验证两部分。",
             "架构侧验证模块是否有用，平台侧验证流程是否跑得通。"),
        ],
        20: [
            ("创新点围绕几何校验、多模型编排、长序列记忆和聚合平台展开。",
             "四个创新点分别回应前面提出的四类问题。"),
            ("各创新点均需通过后续消融和对比实验支撑。",
             "每个创新点都有对应的验证实验。"),
        ],
        23: [
            ("主要风险来自训练质量、算力约束、多模型效果和校验标定复杂度。",
             "四个主要风险，每个都有对应的退路和分阶段应对策略。"),
        ],
    }

    misses = []
    for slide_no, pairs in replacements.items():
        slide = prs.slides[slide_no - 1]
        for old, new in pairs:
            count = replace_exact(slide, old, new)
            if count == 0:
                count = replace_contains(slide, old, new)
            if count == 0:
                misses.append((slide_no, old))

    prs.save(OUT)
    print(OUT)
    if misses:
        print("MISSES:")
        for slide_no, old in misses:
            print(f"Slide {slide_no}: {old}")


if __name__ == "__main__":
    main()
